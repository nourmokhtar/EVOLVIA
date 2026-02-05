import os
import fitz  # PyMuPDF
from pptx import Presentation
import base64
import traceback
from typing import List, Dict
from .agents.deck_analysis_state import SlideData

class DeckService:
    """
    Service to handle pitch deck extraction and multi-agent analysis.
    """
    
    async def extract_deck_data(self, file_path: str, file_type: str) -> List[Dict]:
        """
        Extracts text and (optionally) images from PDF or PPTX.
        """
        slides = []
        try:
            if file_type == "pdf":
                doc = fitz.open(file_path)
                for page_num, page in enumerate(doc):
                    text = page.get_text()
                    # Create a low-res thumbnail for the Design Agent
                    pix = page.get_pixmap(matrix=fitz.Matrix(0.5, 0.5)) 
                    img_base64 = base64.b64encode(pix.tobytes("png")).decode('utf-8')
                    
                    # Structural Metadata
                    word_count = len(text.split())
                    blocks = page.get_text("blocks")
                    
                    slides.append({
                        "page_number": page_num + 1,
                        "text": text,
                        "image_base64": img_base64,
                        "metadata": {
                            "word_count": word_count,
                            "block_count": len(blocks),
                            "is_heavy": word_count > 50 # Heuristic for too much text
                        }
                    })
                doc.close()
            
            elif file_type == "pptx":
                prs = Presentation(file_path)
                for i, slide in enumerate(prs.slides):
                    text_runs = []
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_runs.append(shape.text)
                    slides.append({
                        "page_number": i + 1,
                        "text": "\n".join(text_runs),
                        "image_base64": None
                    })
        except Exception as e:
            print(f"Error extracting deck data: {e}")
            traceback.print_exc()
        
        return slides

    async def analyze_deck(self, file_path: str, file_name: str):
        """
        Main entry point for deck analysis.
        """
        try:
            file_ext = file_name.split('.')[-1].lower()
            if file_ext not in ['pdf', 'pptx']:
                raise ValueError("Unsupported file type. Please upload PDF or PPTX.")

            # 1. Extract Data
            slides = await self.extract_deck_data(file_path, file_ext)
            
            if not slides:
                raise ValueError("Could not extract any content from the deck.")

            # 2. Initial State
            initial_state = {
                "file_name": file_name,
                "file_type": file_ext,
                "slides": slides,
                "presentation_type": "",
                "content_analysis": {},
                "design_analysis": {},
                "strategy_analysis": {},
                "messages": [],
                "overall_score": 0,
                "feedback_summary": "",
                "recommendations": [],
                "presentation_score": 0,
                "content_score": 0,
                "strategy_score": 0
            }

            # 3. Invoke Graph
            from .agents.deck_analysis_graph import deck_graph
            
            # Opik Tracing (Disabled due to S3 upload errors with large frames)
            config = {}
            # try:
            #     from opik.integrations.langchain import OpikTracer
            #     project_name = os.getenv("OPIK_PROJECT_NAME", "evolvia-deck-analysis")
            #     opik_tracer = OpikTracer(project_name=project_name)
            #     config = {"callbacks": [opik_tracer]}
            # except:
            #     pass

            result = await deck_graph.ainvoke(initial_state, config=config)
            
            return {
                "overall_score": result.get("overall_score", 0),
                "content_analysis": result.get("content_analysis", {}),
                "design_analysis": result.get("design_analysis", {}),
                "strategy_analysis": result.get("strategy_analysis", {}),
                "summary": result.get("feedback_summary", ""),
                "recommendations": result.get("recommendations", []),
                "scores": {
                    "content": result.get("content_score", 0),
                    "design": result.get("presentation_score", 0),
                    "strategy": result.get("strategy_score", 0)
                }
            }

        except Exception as e:
            print(f"Critical error in DeckService: {e}")
            traceback.print_exc()
            return {
                "overall_score": 0,
                "summary": f"Error: {str(e)}",
                "recommendations": ["Check file format and try again."]
            }

    async def extract_slides_only(self, file_path: str, file_name: str):
        """Extracts high-quality slide images for presentation mode."""
        try:
            file_ext = file_name.split('.')[-1].lower()
            slides = []
            
            if file_ext == "pdf":
                doc = fitz.open(file_path)
                for page_num, page in enumerate(doc):
                    # Higher res for presentation mode (1.5x)
                    pix = page.get_pixmap(matrix=fitz.Matrix(1.5, 1.5)) 
                    img_base64 = base64.b64encode(pix.tobytes("png")).decode('utf-8')
                    slides.append({
                        "page_number": page_num + 1,
                        "image": img_base64
                    })
                doc.close()
            elif file_ext == "pptx":
                # For PPTX, we'd ideally convert to PDF first or use a specialized library
                # For now, we'll return text or a placeholder if image extraction is too complex
                # But let's try to handle PDF as the primary "presenter" format
                return {"error": "Presentation mode currently supports PDF decks for high-quality rendering."}

            return {"slides": slides}
        except Exception as e:
            print(f"Extraction error: {e}")
            return {"error": str(e)}

deck_service = DeckService()
